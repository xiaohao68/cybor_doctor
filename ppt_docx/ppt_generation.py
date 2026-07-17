'''将大模型生成的json数据转换为ppt，可修改代码自定义ppt的样式'''
import datetime
import hashlib
import os
import time


#qn函数用于创建命名空间的QName对象
from pptx.oxml.ns import qn
from typing import Dict
from env import app_root



ppt_cache_dir = os.path.join(app_root(), "data/cache/ppt")

# 如果文件夹路径不存在，先创建
if not os.path.exists(ppt_cache_dir):
    os.makedirs(ppt_cache_dir)

def make_path(text):
    """生成唯一的文件路径"""
    #haslib.sha256()函数用于生成sha256哈希值，返回一个sha256对象,只能处理字节类型数据
    #hexdigest()方法用于将sha256对象转换为十六进制字符串
    #encode("utf-8")将字符串编码为utf-8字节类型数据 decode("utf-8")将字节类型数据解码为字符串
    file_name = hashlib.sha256(text.encode("utf-8")).hexdigest()  ## 也可以使用uuid
    return os.path.join(ppt_cache_dir, f"{file_name}.pptx")

def render_ppt(ppt_content: Dict) -> str:
    """生成 ppt 文件"""
    #Presentation类用于创建ppt文件，打开已有PPT，添加幻灯片，保存ppt
    from pptx import Presentation
    ppt = Presentation()

    # PPT首页  slides示ppt文件的幻灯片列表
    #add_slide()方法用于添加幻灯片，返回一个Slide对象
    #slide_layouts[0]为title&subtitle layout布局
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    #placeholders属性用于获取幻灯片上的占位符对象列表 
    # text属性用于设置整个占位符纯文本，比如标题，副标题
    #第一个占位符为标题，第二个占位符为副标题
    slide.placeholders[0].text = ppt_content["title"]
    slide.placeholders[1].text = "--来自「赛博华佗」"

    # 内容页
    print(f"总共{len(ppt_content['pages'])}页")
    for i, page in enumerate(ppt_content["pages"]):
        print("生成第%d页:%s" % (i + 1, page["title"]))
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # title&content layout
        
        # 标题  text属性用于设置占位符上的文本 
        slide.placeholders[0].text = page["title"]
        #text_frame属性用于多段落，多格式场景
        # 正文 text_frame功能是用于操作文本框中的文本，包括添加段落、设置段落级别、设置段落文本等
        text_frame = slide.placeholders[1].text_frame  # 获取文本框的text_frame对象
        

        for sub_content in page["content"]:
            print(sub_content)
            
            # 一级正文  add_paragraph()方法用于添加段落，返回一个Paragraph对象
            #sub_title为一级标题段落对象
            #sub_title.text为一级标题文本 
            #sub_title.level为段落级别，2表示二级标题，3表示三级标题，4表示四级标题
            sub_title = text_frame.add_paragraph()
            sub_title.text, sub_title.level = sub_content["title"], 2
            
            # 二级正文  add_paragraph()方法用于添加段落，返回一个Paragraph对象
            #sub_description为二级标题段落对象
            #sub_description.text为二级标题文本 
            #sub_description.level为段落级别，3表示三级标题，4表示四级标题
            sub_description = text_frame.add_paragraph()
            sub_description.text, sub_description.level = sub_content["description"], 3
            
    _output_file = make_path(str(time.time()))
    #save()方法用于保存ppt文件，参数为保存路径，返回值为None
    #_output_file为保存路径，格式为pptx
    ppt.save(_output_file)

    return _output_file
